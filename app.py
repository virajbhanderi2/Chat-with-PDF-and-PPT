import os
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

from flask import Flask, request, jsonify
from flask_cors import CORS
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from rag import load_rag_chain

app = Flask(__name__)
CORS(app)

rag_chain = None


def extract_and_split(file_obj, filename):
    """Extract text from PDF/PPT and split into chunks."""
    text = ""
    if filename.lower().endswith(".pdf"):
        reader = PdfReader(file_obj)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
    elif filename.lower().endswith((".ppt", ".pptx")):
        prs = Presentation(file_obj)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=400, chunk_overlap=150
    )
    return splitter.split_text(text)


@app.route("/")
def home():
    return jsonify({"status": "running"}), 200


@app.route("/api/v1/load", methods=["POST"])
def load_pdf():
    """Upload and process a document file."""
    global rag_chain

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    uploaded_file = request.files["file"]

    allowed_extensions = (".pdf", ".ppt", ".pptx")
    if not uploaded_file.filename.lower().endswith(allowed_extensions):
        return jsonify({"error": "File must be a PDF or PPT"}), 400

    try:
        chunks = extract_and_split(uploaded_file, uploaded_file.filename)
        rag_chain = load_rag_chain(chunks)
        return jsonify({
            "success": True,
            "chunks": len(chunks),
            "filename": uploaded_file.filename
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/v1/ask", methods=["POST"])
def ask():
    """Ask a question about the uploaded PDF."""
    global rag_chain

    if rag_chain is None:
        return jsonify({"error": "No document loaded. Upload a document first."}), 400

    data = request.get_json()
    question = data.get("question", "").strip()

    if not question:
        return jsonify({"error": "Question cannot be empty"}), 400

    try:
        answer = rag_chain.invoke(question)
        return jsonify({"answer": str(answer)}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5000, debug=True)
